"""PowerPoint parser implementation."""

from pathlib import Path

from .base import BaseDocumentParser
from .models import ContentType, ParsedContent


class PptDocumentParser(BaseDocumentParser):
    """Parser for PPT/PPTX documents."""

    def parse(self, file_path: Path) -> list[ParsedContent]:
        """Parse PPT/PPTX and return structured content."""
        path = Path(file_path)
        converted_path: Path | None = None
        try:
            if path.suffix.lower() == ".ppt":
                from .legacy_office import convert_ppt_to_pptx
                converted_path = convert_ppt_to_pptx(path)
                parse_path = converted_path
            else:
                parse_path = path

            from pptx import Presentation  # type: ignore
            prs = Presentation(str(parse_path))

            contents: list[ParsedContent] = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_parts.append(text)

                if slide_parts:
                    contents.append(
                        ParsedContent(
                            content_type=ContentType.TEXT,
                            text=f"[幻灯片 {slide_idx}]\n" + "\n".join(slide_parts),
                            metadata={"slide_number": slide_idx},
                        )
                    )

            return contents
        finally:
            if converted_path is not None and converted_path.exists():
                try:
                    converted_path.unlink()
                except OSError:
                    pass
