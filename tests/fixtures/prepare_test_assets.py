"""
测试素材生成脚本

生成 DOCX、XLSX、PPTX、PNG、JPG 等测试文件到 tests/fixtures/。
运行：在 enterprise_rag 根目录执行
    python tests/fixtures/prepare_test_assets.py

依赖：python-docx, openpyxl, python-pptx（后端已安装）
      Pillow 或标准库（用于 PNG/JPG）
"""

from __future__ import annotations

import base64
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent


def _ensure_dir(p: Path) -> Path:
    p.mkdir(parents=True, exist_ok=True)
    return p


def create_docx(out: Path) -> None:
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        print("跳过 DOCX：未安装 python-docx")
        return

    doc = Document()
    doc.add_paragraph("Enterprise RAG 测试文档 - DOCX")
    doc.add_paragraph("这是用于 T0/O3 测试的 Word 样本文档。")
    doc.add_paragraph("支持格式：TXT, MD, PDF, DOC, DOCX, XLS, XLSX, PPT, PPTX, PNG, JPG。")
    doc.add_paragraph("测试数据：北京、上海、广州、深圳。")
    doc.save(out)
    print(f"已生成: {out}")


def create_xlsx(out: Path) -> None:
    try:
        from openpyxl import Workbook
    except ImportError:
        print("跳过 XLSX：未安装 openpyxl")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "测试页"
    ws["A1"] = "企业 RAG 测试数据"
    ws["A2"] = "城市"
    ws["B2"] = "人口(万)"
    ws["A3"], ws["B3"] = "北京", 2189
    ws["A4"], ws["B4"] = "上海", 2487
    ws["A5"], ws["B5"] = "广州", 1868
    ws["A6"], ws["B6"] = "深圳", 1768
    wb.save(out)
    print(f"已生成: {out}")


def create_pptx(out: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("跳过 PPTX：未安装 python-pptx")
        return

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(2))
    tf = tx.text_frame
    tf.text = "Enterprise RAG 测试 - PPTX"
    p = tf.add_paragraph()
    p.text = "用于 T0/O3 模拟真人测试的 PowerPoint 样本文档。"
    p.font.size = Pt(18)
    prs.save(out)
    print(f"已生成: {out}")


def _create_png_minimal(out: Path) -> None:
    """使用预编码的最小有效 PNG（1x1 灰色像素）。"""
    b64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
    out.write_bytes(base64.b64decode(b64))
    print(f"已生成(最小 PNG): {out}")


def create_png(out: Path) -> None:
    try:
        from PIL import Image

        img = Image.new("RGB", (100, 100), color=(70, 130, 180))
        img.save(out, "PNG")
        print(f"已生成: {out}")
    except ImportError:
        _create_png_minimal(out)


def create_jpg(out: Path) -> None:
    try:
        from PIL import Image

        img = Image.new("RGB", (100, 100), color=(70, 130, 180))
        img.save(out, "JPEG", quality=85)
        print(f"已生成: {out}")
    except ImportError:
        print("跳过 JPG：需要 Pillow，可仅用 PNG 测试")


def main() -> None:
    fixtures = _ensure_dir(SCRIPT_DIR)
    print("生成测试素材到:", fixtures)

    create_docx(fixtures / "test_sample.docx")
    create_xlsx(fixtures / "test_sample.xlsx")
    create_pptx(fixtures / "test_sample.pptx")
    create_png(fixtures / "test_sample.png")
    create_jpg(fixtures / "test_sample.jpg")

    print("完成。可用文件:", sorted(p.name for p in fixtures.iterdir() if p.is_file()))


if __name__ == "__main__":
    main()
