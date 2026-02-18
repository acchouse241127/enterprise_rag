"""Generate test documents for Phase 1.2 parsing tests."""

import sys
from pathlib import Path

FIXTURES_DIR = Path(__file__).resolve().parent.parent / "fixtures"
FIXTURES_DIR.mkdir(parents=True, exist_ok=True)


def gen_docx() -> None:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading("测试文档 - Word", 0)
    doc.add_paragraph("这是 Phase 1.2 的 Word 测试文档。")
    doc.add_paragraph("Enterprise RAG System supports .docx parsing.")
    doc.add_paragraph("包含表格、段落和样式。")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "列1"
    table.rows[0].cells[1].text = "列2"
    table.rows[1].cells[0].text = "A"
    table.rows[1].cells[1].text = "B"
    doc.save(FIXTURES_DIR / "test_normal.docx")
    print("Generated: test_normal.docx")


def gen_xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "测试标题"
    ws["A2"] = "Phase 1.2 Excel 测试"
    ws["B1"] = "数值"
    ws["B2"] = 100
    ws2 = wb.create_sheet("Sheet2")
    ws2["A1"] = "第二页"
    wb.save(FIXTURES_DIR / "test_normal.xlsx")
    print("Generated: test_normal.xlsx")


def gen_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "测试文档 - PPT"
    slide.placeholders[1].text = "Phase 1.2 PowerPoint 测试"
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])
    slide2.shapes.title.text = "第二页"
    slide2.placeholders[1].text = "Enterprise RAG System"
    prs.save(FIXTURES_DIR / "test_normal.pptx")
    print("Generated: test_normal.pptx")


def gen_pdf() -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Test PDF - Phase 1.2", fontsize=24)
    page.insert_text((72, 120), "Enterprise RAG System supports PDF parsing.", fontsize=12)
    page.insert_text((72, 144), "PDF 测试文档。", fontsize=12)
    doc.save(FIXTURES_DIR / "test_normal.pdf")
    doc.close()
    print("Generated: test_normal.pdf")


def main() -> None:
    try:
        gen_docx()
    except Exception as e:
        print(f"Skip docx: {e}")
    try:
        gen_xlsx()
    except Exception as e:
        print(f"Skip xlsx: {e}")
    try:
        gen_pptx()
    except Exception as e:
        print(f"Skip pptx: {e}")
    try:
        gen_pdf()
    except Exception as e:
        print(f"Skip pdf: {e}")


if __name__ == "__main__":
    main()
